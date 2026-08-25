"""生成帮扶村振兴管理系统介绍PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

os.chdir(os.path.join(os.path.dirname(__file__), ".."))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色方案 ──
DARK_GREEN = RGBColor(0x1B, 0x43, 0x32)
MID_GREEN = RGBColor(0x2D, 0x6A, 0x4F)
LIGHT_GREEN = RGBColor(0x40, 0x91, 0x6C)
ACCENT_GOLD = RGBColor(0xE1, 0x9A, 0x25)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x21, 0x25, 0x29)
SUBTITLE_GRAY = RGBColor(0xAD, 0xB5, 0xBD)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, title, content_lines, icon_text="", bg_color=MID_GREEN):
    card = add_shape_bg(slide, left, top, width, height, bg_color)
    card.shadow.inherit = False
    if icon_text:
        add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(0.6), Inches(0.5), icon_text, font_size=28, color=ACCENT_GOLD, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.6), width - Inches(0.6), Inches(0.4), title, font_size=18, color=ACCENT_GOLD, bold=True)
    add_bullet_list(slide, left + Inches(0.3), top + Inches(1.05), width - Inches(0.6), height - Inches(1.2), content_lines, font_size=13, color=WHITE, spacing=Pt(4))
    return card


# ════════════════════════════════════════════
# Slide 1: 封面
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_GREEN)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GOLD)
add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_GOLD)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "帮扶村振兴管理系统", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
             "Assistance Management Information System", font_size=24, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
             "乡村振兴 · 数字化管理", font_size=20, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(4.5), Inches(4.8), Inches(4.333), Inches(0.03), ACCENT_GOLD)

add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "技术架构：Vue 3 + FastAPI + Electron + SQLite", font_size=16, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
             "版本 v1.10.0  |  2026年8月", font_size=14, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 2: 系统概述
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "系统概述", font_size=36, color=WHITE, bold=True)
add_shape_bg(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0),
             "本系统是面向帮扶工作的综合管理平台，实现对帮扶村、帮扶项目、资金管理、学校共建等核心业务的全流程数字化管理。"
             "系统采用单机离线部署模式，支持数据加密导入导出，满足涉密环境下的使用需求。",
             font_size=17, color=WHITE)

cards_data = [
    ("📊", "帮扶村管理", ["10大板块年度数据", "人口/收入/产业/基建", "党建/医疗/教育/就业", "消费/兵力投入/委员会"]),
    ("💰", "资金全生命周期", ["预算编制与审批", "拨付/使用/结算", "异常检测与预警", "绩效评估报告"]),
    ("📋", "项目与审批", ["项目立项/任务/里程碑", "多级审批工作流", "文件管理与变更历史", "进度跟踪与统计"]),
    ("🏫", "学校共建管理", ["学校信息CRUD", "奖学金学生管理", "共建项目管理", "附件与导入导出"]),
]
for i, (icon, title, lines) in enumerate(cards_data):
    col = i % 4
    add_card(slide, Inches(0.8 + col * 3.1), Inches(3.0), Inches(2.8), Inches(3.8), title, lines, icon)

# ════════════════════════════════════════════
# Slide 3: 技术架构
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "技术架构", font_size=36, color=WHITE, bold=True)
add_shape_bg(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), ACCENT_GOLD)

# Architecture layers
layers = [
    ("桌面层", "Electron 28 + NSIS 安装包", "跨平台桌面应用，自动更新，系统托盘", Inches(1.5)),
    ("前端层", "Vue 3 + TypeScript + Element Plus + ECharts", "响应式SPA，Pinia状态管理，Vue Router路由守卫", Inches(2.7)),
    ("后端层", "FastAPI + SQLAlchemy + Pydantic", "RESTful API，JWT认证，RBAC权限，审计日志", Inches(3.9)),
    ("数据层", "SQLite + Alembic 迁移", "单文件数据库，AES-256加密，自动备份", Inches(5.1)),
]
for name, tech, desc, top in layers:
    add_shape_bg(slide, Inches(1.5), top, Inches(10.3), Inches(0.95), MID_GREEN)
    add_text_box(slide, Inches(1.8), top + Inches(0.05), Inches(2), Inches(0.4), name, font_size=20, color=ACCENT_GOLD, bold=True)
    add_text_box(slide, Inches(4.0), top + Inches(0.05), Inches(7.5), Inches(0.4), tech, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(4.0), top + Inches(0.45), Inches(7.5), Inches(0.4), desc, font_size=13, color=SUBTITLE_GRAY)

# Arrows between layers
for top in [Inches(2.45), Inches(3.65), Inches(4.85)]:
    add_text_box(slide, Inches(6.2), top, Inches(1), Inches(0.3), "▼", font_size=16, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 4: 核心功能模块
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "核心功能模块", font_size=36, color=WHITE, bold=True)
add_shape_bg(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), ACCENT_GOLD)

modules = [
    ("🏘️", "帮扶村管理", "10大板块年度数据录入\n委员会/人口/收入/产业\n基建/党建/医疗/教育\n消费/就业/兵力投入"),
    ("💵", "资金管理", "预算编制与执行跟踪\n多级审批与拨付\n异常检测与预警\n绩效评估与结算"),
    ("📁", "项目管理", "项目立项与任务分解\n里程碑进度跟踪\n文件管理与变更历史\n经费关联与统计"),
    ("✅", "审批工作流", "可配置多级审批\n批量审批/转办/撤回\n变更差异对比\n审批提醒与催办"),
    ("📊", "数据分析", "多维度统计看板\n村庄对比与趋势分析\nKPI指标与排名\nAI辅助预测与建议"),
    ("🗺️", "地图可视化", "帮扶村/学校地理标注\n区域分布热力图\n距离计算与路径\n离线地图瓦片"),
    ("🔐", "安全体系", "JWT + 2FA双因素认证\nRBAC细粒度权限\n零信任安全评估\nAES-256数据加密"),
    ("📦", "数据交换", "加密数据包导入导出\n增量同步与冲突解决\n一键报表生成\n数据质量校验"),
]
for i, (icon, title, desc) in enumerate(modules):
    row = i // 4
    col = i % 4
    left = Inches(0.6 + col * 3.15)
    top = Inches(1.5 + row * 2.9)
    add_shape_bg(slide, left, top, Inches(2.9), Inches(2.6), MID_GREEN)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.4), icon, font_size=26, color=ACCENT_GOLD)
    add_text_box(slide, left + Inches(0.7), top + Inches(0.15), Inches(2.0), Inches(0.4), title, font_size=17, color=ACCENT_GOLD, bold=True)
    lines = desc.split('\n')
    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.65), Inches(2.5), Inches(1.8), lines, font_size=12, color=WHITE, spacing=Pt(3))

# ════════════════════════════════════════════
# Slide 5: 安全与部署
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "安全体系与部署方案", font_size=36, color=WHITE, bold=True)
add_shape_bg(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), ACCENT_GOLD)

# Security cards
sec_items = [
    ("🔑", "认证安全", ["JWT双Token机制", "TOTP双因素认证", "机器码绑定验证", "账户锁定与速率限制"]),
    ("🛡️", "权限控制", ["RBAC角色权限", "组织树数据隔离", "菜单级前端权限", "零信任安全评估"]),
    ("🔒", "数据安全", ["AES-256-GCM加密", "CSRF双重提交Cookie", "安全响应头", "审计日志全链路"]),
    ("📡", "部署方案", ["Electron离线单机部署", "NSIS安装包(x64)", "VC++运行时自动安装", "ARM64 Linux DEB包"]),
]
for i, (icon, title, lines) in enumerate(sec_items):
    add_card(slide, Inches(0.6 + i * 3.15), Inches(1.5), Inches(2.9), Inches(3.2), title, lines, icon)

# Deployment flow
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5), "部署流程", font_size=22, color=ACCENT_GOLD, bold=True)
steps = ["下载安装包", "运行NSIS安装程序", "自动安装VC++运行时", "创建桌面快捷方式", "启动系统(自动初始化)"]
for i, step in enumerate(steps):
    left = Inches(0.8 + i * 2.5)
    add_shape_bg(slide, left, Inches(5.8), Inches(2.2), Inches(0.8), LIGHT_GREEN)
    add_text_box(slide, left + Inches(0.1), Inches(5.85), Inches(2.0), Inches(0.7), f"{i+1}. {step}", font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, left + Inches(2.15), Inches(5.9), Inches(0.4), Inches(0.5), "→", font_size=20, color=ACCENT_GOLD)

# ════════════════════════════════════════════
# Slide 6: 项目规模与质量
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "项目规模与质量保障", font_size=36, color=WHITE, bold=True)
add_shape_bg(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), ACCENT_GOLD)

stats = [
    ("699+", "API 端点", "覆盖42个路由模块"),
    ("10,000+", "后端测试", "pytest全量通过"),
    ("1,622", "前端测试", "vitest 125个文件"),
    ("0", "类型错误", "vue-tsc严格模式"),
    ("50+", "Vue 页面", "完整业务功能覆盖"),
    ("7", "CI 工作流", "自动化构建与测试"),
]
for i, (num, label, desc) in enumerate(stats):
    col = i % 3
    row = i // 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.6 + row * 2.5)
    add_shape_bg(slide, left, top, Inches(3.7), Inches(2.1), MID_GREEN)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(3.1), Inches(0.8), num, font_size=44, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.0), Inches(3.1), Inches(0.4), label, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.5), Inches(3.1), Inches(0.4), desc, font_size=13, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 7: 总结
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_GREEN)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GOLD)
add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_GOLD)

add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.0),
             "系统特色与优势", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

features = [
    "🔐  完全离线运行，满足涉密环境要求，数据不出本机",
    "📊  10大板块年度数据全覆盖，资金全生命周期管理",
    "🤖  AI辅助分析预测，智能异常检测与工作建议",
    "🗺️  地图可视化展示，帮扶村与学校地理分布一目了然",
    "⚡  一键数据加密交换，支持多级单位间安全数据传输",
    "🛡️  零信任安全架构，RBAC权限 + 2FA + 审计全链路",
    "📦  Electron桌面应用，NSIS安装包一键部署，自动更新",
]
add_bullet_list(slide, Inches(2.0), Inches(2.5), Inches(9.3), Inches(4.0), features, font_size=19, color=WHITE, spacing=Pt(14))

add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
             "帮扶村振兴管理系统  ·  让帮扶工作更精准、更高效、更安全", font_size=16, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = "deliverables/帮扶村振兴管理系统介绍.pptx"
os.makedirs("deliverables", exist_ok=True)
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
