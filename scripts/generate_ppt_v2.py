"""生成帮扶村振兴管理系统详细介绍PPT（13页精美版）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

os.chdir(os.path.join(os.path.dirname(__file__), ".."))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DG = RGBColor(0x1B, 0x43, 0x32)
MG = RGBColor(0x2D, 0x6A, 0x4F)
LG = RGBColor(0x40, 0x91, 0x6C)
GOLD = RGBColor(0xE1, 0x9A, 0x25)
W = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAD, 0xB5, 0xBD)
DARK = RGBColor(0x21, 0x25, 0x29)
FONT = "微软雅黑"


def bg(slide, color=DG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def bar(slide, top=Inches(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, prs.slide_width, Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()


def txt(slide, l, t, w, h, text, sz=18, c=W, b=False, al=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = FONT
    p.alignment = al
    return tb


def bullets(slide, l, t, w, h, items, sz=15, c=W, sp=Pt(6)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = FONT
        p.space_after = sp
    return tb


def card(slide, l, t, w, h, title, lines, icon="", bgc=MG):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = bgc; s.line.fill.background()
    if icon:
        txt(slide, l + Inches(0.2), t + Inches(0.15), Inches(0.5), Inches(0.4), icon, 24, GOLD, True)
    txt(slide, l + Inches(0.2), t + Inches(0.55), w - Inches(0.4), Inches(0.35), title, 16, GOLD, True)
    bullets(slide, l + Inches(0.2), t + Inches(0.95), w - Inches(0.4), h - Inches(1.1), lines, 12, W, Pt(3))


def title_slide(slide, title, sub=""):
    bg(slide); bar(slide); bar(slide, Inches(7.44))
    txt(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7), title, 34, W, True)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(2), Inches(0.05))
    s.fill.solid(); s.fill.fore_color.rgb = GOLD; s.line.fill.background()
    if sub:
        txt(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4), sub, 15, GRAY)


# ═══ Slide 1: 封面 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s); bar(s); bar(s, Inches(7.44))
txt(s, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2), "帮扶村振兴管理系统", 52, W, True, PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(2.5), Inches(10), Inches(0.7), "Assistance Management Information System", 24, GOLD, False, PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5), "乡村振兴 · 数字化管理 · 离线安全", 20, GRAY, False, PP_ALIGN.CENTER)
s2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.3), Inches(4.333), Inches(0.03))
s2.fill.solid(); s2.fill.fore_color.rgb = GOLD; s2.line.fill.background()
txt(s, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5), "Vue 3 + FastAPI + Electron + SQLite | 完全离线单机部署", 16, GRAY, False, PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5), "版本 v1.10.0  |  2026年8月", 14, GRAY, False, PP_ALIGN.CENTER)

# ═══ Slide 2: 目录 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "目  录")
items = [
    "01  项目背景与建设目标",
    "02  系统概述与适用场景",
    "03  技术架构设计",
    "04  核心功能 — 帮扶村管理",
    "05  核心功能 — 资金全生命周期",
    "06  核心功能 — 项目与审批",
    "07  核心功能 — 数据分析与AI",
    "08  安全体系",
    "09  部署与运维",
    "10  项目质量保障",
    "11  总结与展望",
]
bullets(s, Inches(2), Inches(1.6), Inches(9), Inches(5.5), items, 20, W, Pt(12))

# ═══ Slide 3: 项目背景 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "项目背景与建设目标")
bullets(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
    "📌 战略背景",
    "  • 乡村振兴国家战略",
    "  • 乡村振兴全面推进",
    "  • 部队帮扶工作信息化需求迫切",
    "",
    "🎯 建设目标",
    "  • 帮扶村基础数据全面数字化",
    "  • 帮扶项目全流程闭环管理",
    "  • 帮扶资金全生命周期追踪",
    "  • 帮扶成效多维度分析评估",
], 16, W, Pt(5))
card(s, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5), "核心痛点", [
    "❶ 帮扶数据分散在纸质台账和Excel中",
    "❷ 资金使用缺乏全流程追踪",
    "❸ 帮扶成效难以量化评估",
    "❹ 多级单位间数据交换不安全",
    "❺ 涉密环境无法使用云端系统",
    "",
    "✅ 本系统逐一解决以上痛点",
])

# ═══ Slide 4: 系统概述 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "系统概述与适用场景")
card(s, Inches(0.6), Inches(1.5), Inches(3.8), Inches(2.5), "系统定位", [
    "面向帮扶乡村工作的",
    "综合管理平台",
    "完全离线、单机部署、数据安全",
])
card(s, Inches(4.7), Inches(1.5), Inches(3.8), Inches(2.5), "适用场景", [
    "• 部队帮扶乡村振兴工作管理",
    "• 共建学校管理",
    "• 帮扶资金预算与审计",
    "• 多级单位间加密数据交换",
])
card(s, Inches(8.8), Inches(1.5), Inches(3.8), Inches(2.5), "用户群体", [
    "• 超级管理员：系统全权限",
    "• 单位管理员：本单位数据管理",
    "• 业务经理：审批与数据分析",
    "• 普通用户：数据录入与查看",
])
card(s, Inches(0.6), Inches(4.3), Inches(12), Inches(2.5), "核心功能一览", [
    "🏘️ 帮扶村管理（10大板块年度数据）  💰 资金全生命周期（预算→审批→拨付→使用→结算→审计）  📋 项目管理（任务/里程碑/附件）",
    "🏫 学校共建（奖学金/附件/统计）  ✅ 多级审批工作流  📊 数据分析看板（KPI/趋势/对比）  🗺️ 地图可视化（离线瓦片）",
    "🤖 AI辅助分析（预测/异常检测/建议）  🔐 安全体系（JWT/2FA/RBAC/零信任/加密）  📦 加密数据包交换  📑 报表导出（Excel/PDF）",
])

# ═══ Slide 5: 技术架构 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "技术架构设计")
layers = [
    ("桌面层", "Electron 28 + NSIS安装包", "跨平台桌面应用，系统托盘，自动更新", Inches(1.5)),
    ("前端层", "Vue 3 + TypeScript + Element Plus + ECharts", "响应式SPA，Pinia状态管理，路由懒加载", Inches(2.8)),
    ("后端层", "FastAPI + SQLAlchemy 2.0 + Pydantic", "RESTful API 699端点，JWT认证，RBAC权限", Inches(4.1)),
    ("数据层", "SQLite WAL + Alembic迁移", "单文件数据库，AES-256加密，自动备份", Inches(5.4)),
]
for name, tech, desc, top in layers:
    s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), top, Inches(10.3), Inches(1.0), )
    sh = s.shapes[-1]; sh.fill.solid(); sh.fill.fore_color.rgb = MG; sh.line.fill.background()
    txt(s, Inches(1.8), top + Inches(0.05), Inches(2), Inches(0.4), name, 20, GOLD, True)
    txt(s, Inches(4.0), top + Inches(0.05), Inches(7.5), Inches(0.4), tech, 16, W, True)
    txt(s, Inches(4.0), top + Inches(0.5), Inches(7.5), Inches(0.4), desc, 13, GRAY)
for top in [Inches(2.5), Inches(3.8), Inches(5.1)]:
    txt(s, Inches(6.2), top, Inches(1), Inches(0.3), "▼", 16, GOLD, False, PP_ALIGN.CENTER)

# ═══ Slide 6: 帮扶村管理 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "核心功能 — 帮扶村管理")
sections = [
    ("👥", "人口数据", ["总户数/总人数", "常住人口/劳动力"]),
    ("💵", "收入数据", ["人均纯收入", "集体经济收入"]),
    ("🏭", "产业帮扶", ["产业类型/投入", "带动就业人数"]),
    ("🏗️", "基础设施", ["道路/水利/电力", "住房改造"]),
    ("🚩", "党建引领", ["党员人数", "组织建设"]),
    ("🏥", "医疗帮扶", ["医疗投入", "健康档案"]),
    ("🎓", "教育帮扶", ["教育投入", "资助学生"]),
    ("💼", "就业帮扶", ["就业培训", "转移就业"]),
    ("🛒", "消费帮扶", ["采购金额", "销售渠道"]),
    ("🎖️", "兵力投入", ["派驻人员", "帮扶天数"]),
]
for i, (icon, title, lines) in enumerate(sections):
    col = i % 5; row = i // 5
    card(s, Inches(0.5 + col * 2.5), Inches(1.5 + row * 2.8), Inches(2.3), Inches(2.5), title, lines, icon)

# ═══ Slide 7: 资金管理 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "核心功能 — 资金全生命周期")
steps = ["预算编制", "审批通过", "资金拨付", "使用跟踪", "结算审核", "绩效评估"]
for i, step in enumerate(steps):
    left = Inches(0.8 + i * 2.1)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(1.8), Inches(0.8))
    sh.fill.solid(); sh.fill.fore_color.rgb = LG if i % 2 == 0 else MG; sh.line.fill.background()
    txt(s, left + Inches(0.1), Inches(1.9), Inches(1.6), Inches(0.6), f"{i+1}. {step}", 15, W, True, PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txt(s, left + Inches(1.75), Inches(1.9), Inches(0.4), Inches(0.5), "→", 20, GOLD, False, PP_ALIGN.CENTER)
card(s, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.5), "预算管理", [
    "• 年度预算编制与分类管理",
    "• 预算执行率实时仪表盘",
    "• 预算预警（超支/即将超支）",
    "• 预算交易明细记录",
    "• 各分类进度条可视化",
])
card(s, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.5), "多维分析", [
    "• 按类型/年度/地区多维统计",
    "• 年度资金趋势面积图",
    "• 利用率趋势曲线",
    "• 年度对比分析",
    "• 资金流向追踪",
])

# ═══ Slide 8: 项目与审批 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "核心功能 — 项目与审批")
card(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5), "帮扶项目管理", [
    "📋 项目立项与基本信息",
    "📝 任务分解与进度跟踪",
    "🎯 里程碑管理与完成确认",
    "💰 经费关联与统计",
    "📎 附件上传/下载/预览",
    "📜 变更历史完整记录",
    "📊 任务完成率进度条",
    "📥 Excel批量导入导出",
])
card(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5), "多级审批工作流", [
    "✅ 可配置多级审批流程",
    "📤 提交审批（自动通知审批人）",
    "👍 通过 / 👎 驳回 / 🔄 转办",
    "↩️ 申请人可撤回",
    "📦 批量审批（管理员）",
    "📊 变更差异对比",
    "⏰ 审批提醒与催办",
    "📜 审批历史完整记录",
])

# ═══ Slide 9: 数据分析与AI ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "核心功能 — 数据分析与AI")
card(s, Inches(0.6), Inches(1.5), Inches(3.8), Inches(2.5), "📊 数据看板", [
    "KPI统计卡片（自动刷新）",
    "年度数据趋势图",
    "村庄对比分析",
    "绩效排名",
])
card(s, Inches(4.7), Inches(1.5), Inches(3.8), Inches(2.5), "🗺️ 地图可视化", [
    "帮扶村/学校地理标注",
    "区域分布展示",
    "离线地图瓦片",
    "坐标编辑与距离计算",
])
card(s, Inches(8.8), Inches(1.5), Inches(3.8), Inches(2.5), "🤖 AI辅助", [
    "收入预测（时间序列）",
    "异常数据检测",
    "智能工作建议",
    "自然语言查询",
])
card(s, Inches(0.6), Inches(4.3), Inches(12), Inches(2.5), "📑 报表与导出", [
    "综合报表生成（年度/季度）  •  Excel/PDF/WPS多格式导出  •  加密数据包（AES-256）多级单位安全交换  •  数据质量校验（完整性/一致性/准确性）",
    "报表订阅（定期自动生成）  •  数据同步（增量/全量）  •  冲突解决机制  •  导入预览与校验",
])

# ═══ Slide 10: 安全体系 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "安全体系")
cards = [
    ("🔑", "认证安全", ["JWT双Token机制", "TOTP双因素认证", "机器码绑定验证", "账户锁定+速率限制"]),
    ("🛡️", "权限控制", ["RBAC角色权限", "组织树数据隔离", "菜单级前端权限", "工作台按角色过滤"]),
    ("🔒", "数据安全", ["AES-256-GCM加密", "CSRF双重提交Cookie", "安全响应头", "审计日志全链路"]),
    ("📡", "离线部署", ["数据不出本机", "加密数据包交换", "NSIS安装包", "VC++运行时内置"]),
]
for i, (icon, title, lines) in enumerate(cards):
    card(s, Inches(0.5 + i * 3.15), Inches(1.5), Inches(2.9), Inches(3.2), title, lines, icon)
card(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), "零信任安全架构", [
    "设备指纹识别 → 动态权限评估 → 最小权限原则 → 持续验证 → 安全事件审计  |  评估维度：设备可信度 + 用户行为 + 网络环境 + 数据敏感度",
])

# ═══ Slide 11: 部署与运维 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "部署与运维")
steps = ["下载安装包", "运行安装程序", "自动安装VC++", "创建桌面快捷方式", "启动系统"]
for i, step in enumerate(steps):
    left = Inches(0.6 + i * 2.5)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(2.2), Inches(0.8))
    sh.fill.solid(); sh.fill.fore_color.rgb = LG; sh.line.fill.background()
    txt(s, left + Inches(0.1), Inches(1.6), Inches(2.0), Inches(0.6), f"{i+1}. {step}", 14, W, True, PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txt(s, left + Inches(2.15), Inches(1.6), Inches(0.4), Inches(0.5), "→", 18, GOLD)
card(s, Inches(0.6), Inches(2.8), Inches(5.8), Inches(3.8), "运维功能", [
    "🔄 自动备份（每日/每周/每月）",
    "🏥 一键系统体检（6维度健康检查）",
    "📊 系统监控仪表板（CPU/内存/磁盘）",
    "📋 审计日志管理",
    "🔧 缓存管理与清理",
    "📝 更新日志与版本管理",
    "🌐 运行环境检查",
])
card(s, Inches(6.8), Inches(2.8), Inches(5.8), Inches(3.8), "系统管理", [
    "👤 用户管理（创建/禁用/重置密码）",
    "🔐 权限配置（角色/菜单/数据范围）",
    "⚙️ 系统配置（配置包导入导出）",
    "🔑 密钥管理（轮换/撤销）",
    "📧 邮件服务配置",
    "🗺️ 地图瓦片管理",
    "🔒 加密设置（数据库加密）",
])

# ═══ Slide 12: 项目质量 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
title_slide(s, "项目质量保障")
stats = [
    ("699+", "API 端点", "42个路由模块"),
    ("10,056", "后端测试", "pytest全量通过"),
    ("1,622", "前端测试", "vitest 125文件"),
    ("0", "类型错误", "vue-tsc严格模式"),
    ("7", "CI 工作流", "自动化构建测试"),
    ("15", "帮助文档", "覆盖全部模块"),
]
for i, (num, label, desc) in enumerate(stats):
    col = i % 3; row = i // 3
    left = Inches(0.8 + col * 4.1); top = Inches(1.5 + row * 2.7)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.3))
    sh.fill.solid(); sh.fill.fore_color.rgb = MG; sh.line.fill.background()
    txt(s, left + Inches(0.3), top + Inches(0.2), Inches(3.1), Inches(0.8), num, 44, GOLD, True, PP_ALIGN.CENTER)
    txt(s, left + Inches(0.3), top + Inches(1.0), Inches(3.1), Inches(0.4), label, 20, W, True, PP_ALIGN.CENTER)
    txt(s, left + Inches(0.3), top + Inches(1.5), Inches(3.1), Inches(0.4), desc, 13, GRAY, False, PP_ALIGN.CENTER)

# ═══ Slide 13: 总结 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s); bar(s); bar(s, Inches(7.44))
txt(s, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8), "总结与展望", 40, W, True, PP_ALIGN.CENTER)
features = [
    "🔐  完全离线运行，满足涉密环境要求，数据不出本机",
    "📊  10大板块年度数据全覆盖，资金全生命周期管理",
    "🤖  AI辅助分析预测，智能异常检测与工作建议",
    "🗺️  地图可视化展示，帮扶村与学校地理分布一目了然",
    "⚡  一键数据加密交换，支持多级单位间安全数据传输",
    "🛡️  零信任安全架构，RBAC权限 + 2FA + 审计全链路",
    "📦  Electron桌面应用，NSIS安装包一键部署",
    "📖  15篇帮助文档，小白也能快速上手",
]
bullets(s, Inches(2.0), Inches(2.0), Inches(9.3), Inches(4.0), features, 19, W, Pt(14))
txt(s, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
    "帮扶村振兴管理系统  ·  让帮扶工作更精准、更高效、更安全", 16, GOLD, False, PP_ALIGN.CENTER)

# ── 保存 ──
out = "deliverables/帮扶村振兴管理系统详细介绍.pptx"
os.makedirs("deliverables", exist_ok=True)
prs.save(out)
print(f"PPT saved: {out} ({len(prs.slides)} slides)")
